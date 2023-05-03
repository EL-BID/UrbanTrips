import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor # ColorFormat, 
from PIL import Image, ImageDraw


def crop_imagen(filePath, reduce=1, altura_max=0, ancho_max=0, save=True, crop_left = 0, crop_top = 0, crop_right = 0, crop_bottom = 0):
    
    # Trim all png images with white background in a folder
    # Usage "python PNGWhiteTrim.py ../someFolder padding"

    image=Image.open(filePath)
    image.load()
    imageSize = image.size #tuple
    
    ## QUITA ESPACIOS EN BLANCO ALREDEDOR DE LA IMAGEN
    # remove alpha channel
    invert_im = image.convert("RGB")
    # invert image (so that white is 0)
    invert_im = ImageOps.invert(invert_im)
    imageBox = invert_im.getbbox()
    cropped=image.crop(imageBox)
    ## FIN DE QUITA ESPACIOS EN BLANCO ALREDEDOR DE LA IMAGEN
    
    #REDUCE TAMAÑO
    _size=[]
    # calculates percentage to reduce image by maintaining proportion
    if altura_max>0: _size.append((altura_max/(cropped.height/38)))
    if ancho_max>0: _size.append((ancho_max/(cropped.width/38)))
    if len(_size) > 0: reduce = min(_size)
    
    if reduce < 1:
        basewidth = int(cropped.width * reduce)
        wpercent = (basewidth/float(cropped.size[0]))
        hsize = int((float(cropped.size[1])*float(wpercent)))
        # cropped.resize actually does the resizing
        cropped = cropped.resize((basewidth,hsize), Image.ANTIALIAS)
    
    if crop_left + crop_top + crop_right + crop_bottom > 0:
        width, height = cropped.size 
        crop_right = width - crop_right 
        crop_bottom = height - crop_bottom            
        cropped=cropped.crop((crop_left, crop_top, crop_right, crop_bottom))

    # save the image as cropped
    if save:
        filePath = filePath[0: filePath.find('.')]+'_cropped'+filePath[filePath.find('.'):len(filePath)]
        cropped.save(filePath)
        return filePath
    else:
        return cropped

def pptx_addtitle(prs, slide='', title='', top=0, left=0, width=10, height=1, new=True, fontsize=24, fontcolor='blue', bold=True):

    blank_slide_layout = prs.slide_layouts[6] # Using layout 6 (blank layout)
    # if new create blank slide
    if new:
        slide = prs.slides.add_slide(blank_slide_layout)

    # Set the slides background colour
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(212, 218, 220) # RGBColor(212, 218, 220) is the color of water on the contextily tiles

    # translates from cm to inches
    top = Inches(top)
    left = Inches(left)
    width = Inches(width)
    height = Inches(height)
    
    # adds a text box onto the slide object
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.auto_size = False
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = title
    p.font.name = 'Gill Sans'
    p.font.color.rgb = RGBColor(64,64,64) # (105,105,105) CSS Dim Grey
    if bold is True:
        p.font.bold = True
    p.font.size = Pt(fontsize)
    p.alignment = PP_ALIGN.CENTER
    
    #p.font.color = fontcolor
    # many more parameters available

    return slide

def pptx_addtext(prs, slide='', text='', top= 0, left=0, width=10, height=1):
    blank_slide_layout = prs.slide_layouts[6]
    
    top = Inches(top)
    left = Inches(left)
    width = Inches(width)
    height = Inches(height)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.title = text
    # p.font.bold = True
    
    p.alignment = PP_ALIGN.RIGHT
    
    return slide

def pptx_addpic(prs, slide, img_path,  left=0, top=0, width=0, altura_max=0, ancho_max=0, crop_left = 0, crop_top = 0, crop_right = 0, crop_bottom = 0):
    # for adding all maps and graphs
    # altura_max and ancho_max in cm
    blank_slide_layout = prs.slide_layouts[6]

    img_path = str(img_path)

    if os.path.exists(img_path):
        # crop_imagen crops the image
        # NB commented out 20200514
        img_path = crop_imagen(img_path, reduce=1, altura_max=altura_max, ancho_max=ancho_max, save=True, crop_left=crop_left, crop_top=crop_top, crop_right=crop_right, crop_bottom=crop_bottom)
        
        # control position
        left = Inches(left)
        top = Inches(top)
        width  = Inches(width)
        # add to the slide
        if width!=0:
            slide_return = slide.shapes.add_picture(img_path, left, top, width) 
        else:
            slide_return = slide.shapes.add_picture(img_path, left, top) 
        
        os.remove(img_path)
        
        return slide_return