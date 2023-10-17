import pandas as pd
import os
from urbantrips.utils.utils import (
    leer_configs_generales,
    traigo_db_path,
    iniciar_conexion_db,
    leer_alias)

import requests

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor  # ColorFormat,
from PIL import Image, ImageDraw, ImageOps


def crop_imagen(filePath, reduce=1, altura_max=0, ancho_max=0, save=True, crop_left=0, crop_top=0, crop_right=0, crop_bottom=0):

    # Trim all png images with white background in a folder
    # Usage "python PNGWhiteTrim.py ../someFolder padding"

    image = Image.open(filePath)
    image.load()
    imageSize = image.size  # tuple

    # QUITA ESPACIOS EN BLANCO ALREDEDOR DE LA IMAGEN
    # remove alpha channel
    invert_im = image.convert("RGB")
    # invert image (so that white is 0)
    invert_im = ImageOps.invert(invert_im)
    imageBox = invert_im.getbbox()
    cropped = image.crop(imageBox)
    # FIN DE QUITA ESPACIOS EN BLANCO ALREDEDOR DE LA IMAGEN

    # REDUCE TAMAÑO
    _size = []
    # calculates percentage to reduce image by maintaining proportion
    if altura_max > 0:
        _size.append((altura_max/(cropped.height/38)))
    if ancho_max > 0:
        _size.append((ancho_max/(cropped.width/38)))
    if len(_size) > 0:
        reduce = min(_size)

    if reduce < 1:
        basewidth = int(cropped.width * reduce)
        wpercent = (basewidth/float(cropped.size[0]))
        hsize = int((float(cropped.size[1])*float(wpercent)))
        # cropped.resize actually does the resizing
        cropped = cropped.resize((basewidth, hsize), Image.ANTIALIAS)

    if crop_left + crop_top + crop_right + crop_bottom > 0:
        width, height = cropped.size
        crop_right = width - crop_right
        crop_bottom = height - crop_bottom
        cropped = cropped.crop((crop_left, crop_top, crop_right, crop_bottom))

    # save the image as cropped
    if save:
        filePath = filePath[0: filePath.find(
            '.')]+'_cropped'+filePath[filePath.find('.'):len(filePath)]
        cropped.save(filePath)
        return filePath
    else:
        return cropped


def pptx_addtitle(prs, slide='', title='', top=0, left=0, width=10, height=1, new=True, fontsize=24, fontcolor='blue', bold=True):

    blank_slide_layout = prs.slide_layouts[6]  # Using layout 6 (blank layout)
    # if new create blank slide
    if new:
        slide = prs.slides.add_slide(blank_slide_layout)

    # # Set the slides background colour
    # background = slide.background
    # fill = background.fill
    # fill.solid()
    # fill.fore_color.rgb = RGBColor(212, 218, 220) # RGBColor(212, 218, 220) is the color of water on the contextily tiles

    # translates from cm to inches
    top = Inches(top)
    left = Inches(left)
    width = Inches(width)
    height = Inches(height)

    # adds a text box onto the slide object
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.auto_size = False
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = title
    p.font.name = 'Gill Sans'
    p.font.color.rgb = RGBColor(64, 64, 64)  # (105,105,105) CSS Dim Grey
    if bold is True:
        p.font.bold = True

    p.font.size = Pt(fontsize)
    p.alignment = PP_ALIGN.CENTER

    # p.font.color = fontcolor
    # many more parameters available

    return slide


def pptx_text(prs, slide='', title='', top=0, left=0, width=10, height=1, fontsize=24, fontcolor='blue', bold=False):

    blank_slide_layout = prs.slide_layouts[6]  # Using layout 6 (blank layout)

    # # Set the slides background colour
    # background = slide.background
    # fill = background.fill
    # fill.solid()
    # fill.fore_color.rgb = RGBColor(212, 218, 220) # RGBColor(212, 218, 220) is the color of water on the contextily tiles

    # translates from cm to inches
    top = Inches(top)
    left = Inches(left)
    width = Inches(width)
    height = Inches(height)

    # adds a text box onto the slide object
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.auto_size = False
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = title
    p.font.name = 'Gill Sans'
    p.font.color.rgb = RGBColor(64, 64, 64)  # (105,105,105) CSS Dim Grey
    if bold is True:
        p.font.bold = True

    p.font.size = Pt(fontsize)
    p.alignment = PP_ALIGN.LEFT

    # p.font.color = fontcolor
    # many more parameters available

    return slide


def pptx_addpic(prs, slide, img_path,  left=0, top=0, width=0, altura_max=0, ancho_max=0, crop_left=0, crop_top=0, crop_right=0, crop_bottom=0):
    # for adding all maps and graphs
    # altura_max and ancho_max in cm
    # blank_slide_layout = prs.slide_layouts[6]

    img_path = str(img_path)

    if os.path.exists(img_path):
        # crop_imagen crops the image
        # NB commented out 20200514
        img_path = crop_imagen(img_path, reduce=1, altura_max=altura_max, ancho_max=ancho_max, save=True,
                               crop_left=crop_left, crop_top=crop_top, crop_right=crop_right, crop_bottom=crop_bottom)

        # control position
        left = Inches(left)
        top = Inches(top)
        width = Inches(width)
        # add to the slide
        if width != 0:
            slide_return = slide.shapes.add_picture(img_path, left, top, width)
        else:
            slide_return = slide.shapes.add_picture(img_path, left, top)

        os.remove(img_path)

        return slide_return


def get_new_slide(prs, desc_dia_titulo):
    # Logo Urbantrips
    try:
        db_path = os.path.join("docs")
        os.makedirs(db_path, exist_ok=True)

        file_logo = os.path.join(
            "docs", "urbantrips_logo.jpg")
        if not os.path.isfile(file_logo):
            # URL of the image file on Github
            url = 'https://raw.githubusercontent.com/EL-BID/UrbanTrips/main/docs/urbantrips_logo.jpg'
            # Send a request to get the content of the image file
            response = requests.get(url)

            # Save the content to a local file
            with open(file_logo, 'wb') as f:
                f.write(response.content)

        slide = pptx_addtitle(prs=prs, slide='',  title='',
                              left=0, top=0, width=24, new=True, fontsize=48)
        pptx_addpic(prs=prs, slide=slide, img_path=file_logo,
                    left=16, top=12.3, width=8)
    except:
        pass

    slide = pptx_addtitle(prs=prs, slide=slide,  title='Urbantrips',
                          left=0, top=0, width=24, new=False, fontsize=48)
    slide = pptx_addtitle(prs=prs, slide=slide,  title=desc_dia_titulo,
                          left=0, top=1, width=24, new=False, fontsize=38)

    return slide


def format_num(num, lpad=10):
    fnum = '{:,}'.format(num).replace(
        ".", "*").replace(",", ".").replace("*", ",")
    if lpad > 0:
        fnum = fnum.rjust(lpad, ' ')
    return fnum


def slide_1(prs,
            indicadores,
            desc_dia,
            tipo_dia,
            desc_dia_titulo):

    df_indicadores = pd.DataFrame([])

    slide = get_new_slide(prs, desc_dia_titulo)

    top_i = 3
    left_i = 1

    slide = pptx_text(prs=prs, slide=slide,  title=f'Información del dataset original',
                      left=left_i, top=2.3, width=18, fontsize=24, bold=True)

    ind_name = 'Cantidad de transacciones totales'
    ind = format_num(
        indicadores.loc[indicadores.detalle == ind_name].indicador.astype(int).values[0])
    slide = pptx_text(prs=prs, slide=slide,
                      title=f'{ind_name}:', left=left_i, top=top_i, width=18, fontsize=18, bold=True)
    slide = pptx_text(prs=prs, slide=slide,  title=ind,
                      left=left_i+4.3, top=top_i, width=18, fontsize=18)

    df_indicadores = pd.concat([df_indicadores, pd.DataFrame([['Información del dataset original', 1, ind_name, ind]], columns=[
                               'Titulo', 'orden', 'Indicador', 'Valor'])], ignore_index=True)

    top_i += 0.5
    ind_name = 'Cantidad de tarjetas únicas'
    ind = format_num(
        indicadores.loc[indicadores.detalle == ind_name].indicador.astype(int).values[0])
    slide = pptx_text(prs=prs, slide=slide,
                      title=f'{ind_name}:', left=left_i, top=top_i, width=18, fontsize=18, bold=True)
    slide = pptx_text(prs=prs, slide=slide,  title=ind,
                      left=left_i+4.3, top=top_i, width=18, fontsize=18)
    df_indicadores = pd.concat([df_indicadores, pd.DataFrame([['Información del dataset original', 1, ind_name, ind]], columns=[
                               'Titulo', 'orden', 'Indicador', 'Valor'])], ignore_index=True)

    top_i += 0.5
    ind_name = 'Cantidad de transacciones limpias'
    ind = format_num(
        indicadores.loc[indicadores.detalle == ind_name].indicador.astype(int).values[0])
    slide = pptx_text(prs=prs, slide=slide,
                      title=f'{ind_name}:', left=left_i, top=top_i, width=18, fontsize=18, bold=True)
    slide = pptx_text(prs=prs, slide=slide,  title=ind,
                      left=left_i+4.3, top=top_i, width=18, fontsize=18)
    df_indicadores = pd.concat([df_indicadores, pd.DataFrame([['Información del dataset original', 1, ind_name, ind]], columns=[
                               'Titulo', 'orden', 'Indicador', 'Valor'])], ignore_index=True)

    top_i += 0.5
    ind_name = 'Cantidad de etapas con destinos validados'
    ind = format_num(
        indicadores.loc[indicadores.detalle == ind_name].indicador.astype(int).values[0])
    ind += ' ('+format_num(indicadores.loc[indicadores.detalle ==
                                           ind_name].porcentaje.values[0], 0)+'%)'
    slide = pptx_text(prs=prs, slide=slide,  title=f'Transacciones válidas \n(Etapas con destinos validados):',
                      left=left_i, top=top_i, width=18, fontsize=18, bold=True)
    slide = pptx_text(prs=prs, slide=slide,  title=ind,
                      left=left_i+4.3, top=top_i, width=18, fontsize=18)
    df_indicadores = pd.concat([df_indicadores, pd.DataFrame([['Información del dataset original', 1, 'Transacciones válidas \n(Etapas con destinos validados):', ind]], columns=[
                               'Titulo', 'orden', 'Indicador', 'Valor'])], ignore_index=True)

    top_i = 3
    left_i = 9

    slide = pptx_text(prs=prs, slide=slide,  title=f'Información procesada',
                      left=left_i, top=2.3, width=18, fontsize=24, bold=True)

    ind_name = 'Cantidad total de viajes expandidos'
    ind = format_num(
        indicadores.loc[indicadores.detalle == ind_name].indicador.astype(int).values[0])
    slide = pptx_text(prs=prs, slide=slide,  title=f'Viajes:',
                      left=left_i, top=top_i, width=18, fontsize=18, bold=True)
    slide = pptx_text(prs=prs, slide=slide,  title=ind,
                      left=left_i+4.3, top=top_i, width=18, fontsize=18)
    df_indicadores = pd.concat([df_indicadores, pd.DataFrame([['Información procesada', 2, 'Viajes', ind]], columns=[
                               'Titulo', 'orden', 'Indicador', 'Valor'])], ignore_index=True)

    top_i += 0.5
    ind_name = 'Cantidad total de etapas'
    ind = format_num(
        indicadores.loc[indicadores.detalle == ind_name].indicador.astype(int).values[0])
    slide = pptx_text(prs=prs, slide=slide,  title=f'Etapas:',
                      left=left_i, top=top_i, width=18, fontsize=18, bold=True)
    slide = pptx_text(prs=prs, slide=slide,  title=ind,
                      left=left_i+4.3, top=top_i, width=18, fontsize=18)
    df_indicadores = pd.concat([df_indicadores, pd.DataFrame([['Información procesada', 2, 'Etapas', ind]], columns=[
                               'Titulo', 'orden', 'Indicador', 'Valor'])], ignore_index=True)

    top_i += 0.5
    ind_name = 'Cantidad total de usuarios'
    ind = format_num(
        indicadores.loc[indicadores.detalle == ind_name].indicador.astype(int).values[0])
    slide = pptx_text(prs=prs, slide=slide,  title=f'Usuarios:',
                      left=left_i, top=top_i, width=18, fontsize=18, bold=True)
    slide = pptx_text(prs=prs, slide=slide,  title=ind,
                      left=left_i+4.3, top=top_i, width=18, fontsize=18)
    df_indicadores = pd.concat([df_indicadores, pd.DataFrame([['Información procesada', 2, 'Usuarios', ind]], columns=[
                               'Titulo', 'orden', 'Indicador', 'Valor'])], ignore_index=True)

    top_i += 0.5
    ind_name = 'Cantidad de viajes cortos (<5kms)'
    ind = format_num(
        indicadores.loc[indicadores.detalle == ind_name].indicador.astype(int).values[0])
    ind += ' ('+format_num(indicadores.loc[indicadores.detalle ==
                                           ind_name].porcentaje.values[0], 0)+'%)'
    slide = pptx_text(prs=prs, slide=slide,  title=f'Viajes cortos (<5kms):',
                      left=left_i, top=top_i, width=18, fontsize=18, bold=True)
    slide = pptx_text(prs=prs, slide=slide,  title=ind,
                      left=left_i+4.3, top=top_i, width=18, fontsize=18)
    df_indicadores = pd.concat([df_indicadores, pd.DataFrame([['Información procesada', 2, 'Viajes cortos (<5kms)', ind]], columns=[
                               'Titulo', 'orden', 'Indicador', 'Valor'])], ignore_index=True)

    top_i += 0.5
    ind_name = 'Cantidad de viajes con transferencia'
    ind = format_num(
        indicadores.loc[indicadores.detalle == ind_name].indicador.astype(int).values[0])
    ind += ' ('+format_num(indicadores.loc[indicadores.detalle ==
                                           ind_name].porcentaje.values[0], 0)+'%)'
    slide = pptx_text(prs=prs, slide=slide,  title=f'Viajes con transferencia:',
                      left=left_i, top=top_i, width=18, fontsize=18, bold=True)
    slide = pptx_text(prs=prs, slide=slide,  title=ind,
                      left=left_i+4.3, top=top_i, width=18, fontsize=18)
    df_indicadores = pd.concat([df_indicadores, pd.DataFrame([['Información procesada', 2, 'Viajes con transferencia', ind]], columns=[
                               'Titulo', 'orden', 'Indicador', 'Valor'])], ignore_index=True)

    top_i += 0.5
    ind_name = 'Distancia de los viajes (promedio en kms)'
    ind = format_num(
        indicadores.loc[indicadores.detalle == ind_name].indicador.astype(int).values[0])
    slide = pptx_text(prs=prs, slide=slide,  title=f'Distancia de viajes (promedio en kms):',
                      left=left_i, top=top_i, width=18, fontsize=18, bold=True)
    slide = pptx_text(prs=prs, slide=slide,  title=ind,
                      left=left_i+4.3, top=top_i, width=18, fontsize=18)
    df_indicadores = pd.concat([df_indicadores, pd.DataFrame([['Información procesada', 2, 'Distancia de los viajes (promedio en kms)', ind]], columns=[
                               'Titulo', 'orden', 'Indicador', 'Valor'])], ignore_index=True)

    top_i += 0.5
    ind_name = 'Distancia de los viajes (mediana en kms)'
    ind = format_num(
        indicadores.loc[indicadores.detalle == ind_name].indicador.astype(int).values[0])
    slide = pptx_text(prs=prs, slide=slide,  title=f'Distancia de viajes (mediana en kms):',
                      left=left_i, top=top_i, width=18, fontsize=18, bold=True)
    slide = pptx_text(prs=prs, slide=slide,  title=ind,
                      left=left_i+4.3, top=top_i, width=18, fontsize=18)
    df_indicadores = pd.concat([df_indicadores, pd.DataFrame([['Información procesada', 2, 'Distancia de viajes (mediana en kms)', ind]], columns=[
                               'Titulo', 'orden', 'Indicador', 'Valor'])], ignore_index=True)

    top_i = 3
    left_i = 17

    slide = pptx_text(prs=prs, slide=slide,  title=f'Partición Modal',
                      left=left_i, top=2.3, width=18, fontsize=24, bold=True)

    for _, i in indicadores[(indicadores.tabla == 'modos viajes') & (indicadores.detalle != 'Cantidad total de viajes expandidos')].iterrows():

        ind_name = i.detalle
        ind = format_num(
            indicadores.loc[indicadores.detalle == ind_name].indicador.astype(int).values[0])
        ind += ' ('+format_num(
            indicadores.loc[indicadores.detalle == ind_name].porcentaje.values[0], 0)+'%)'
        slide = pptx_text(prs=prs, slide=slide,  title=i.detalle,
                          left=left_i, top=top_i, width=18, fontsize=18, bold=True)
        slide = pptx_text(prs=prs, slide=slide,  title=ind,
                          left=left_i+4.3, top=top_i, width=18, fontsize=18)
        top_i += 0.5

        df_indicadores = pd.concat([df_indicadores, pd.DataFrame([['Partición modal', 3, ind_name, ind]], columns=[
                                   'Titulo', 'orden', 'Indicador', 'Valor'])], ignore_index=True)

    df_indicadores['desc_dia'] = desc_dia
    df_indicadores['tipo_dia'] = tipo_dia

    # conn_dash = iniciar_conexion_db(tipo='dash')

    # df_indicadores_ant = pd.read_sql_query(
    # """
    # SELECT *
    # FROM indicadores
    # """,
    # conn_dash,
    # )

    # df_indicadores_ant = df_indicadores_ant[~(
    #            (df_indicadores_ant.desc_dia==desc_dia)&
    #            (df_indicadores_ant.tipo_dia==tipo_dia)
    #           )]

    # df_indicadores=pd.concat([df_indicadores_ant, df_indicadores], ignore_index=True)

    # df_indicadores.to_sql("indicadores", conn_dash, if_exists="replace", index=False)
    # conn_dash.close()

    return prs


def create_ppt():

    print('')
    print('create_ppt')
    print('----------')

    meses = ['Enero', 'Febrero', 'Marzo', 'Abril', 'Mayo', 'Junio',
             'Julio', 'Agosto', 'Septiembre', 'Octubre', 'Noviembre',
             'Diciembre']

    pd.options.mode.chained_assignment = None

    alias = leer_alias()

    configs = leer_configs_generales()

    try:
        zonificaciones = configs['zonificaciones']
    except KeyError:
        zonificaciones = []

    geo_files = []

    if zonificaciones:
        for n in range(0, 5):

            try:
                file_zona = zonificaciones[f"geo{n+1}"]
                var_zona = zonificaciones[f"var{n+1}"]
                geo_files += [[file_zona, var_zona]]

            except KeyError:
                pass

    geo_files += [['Zona_voi.geojson', 'Zona_voi']]

    # Leer informacion de viajes y distancias
    conn_data = iniciar_conexion_db(tipo='data')
    conn_insumos = iniciar_conexion_db(tipo='insumos')
    conn_dash = iniciar_conexion_db(tipo='dash')

    viajes = pd.read_sql_query(
        """
        SELECT *
        FROM viajes
        """,
        conn_data,
    )

    indicadores = pd.read_sql_query(
        """
        SELECT *
        FROM indicadores
        """,
        conn_data,
    )

    indicadores_dash = pd.read_sql_query(
        """
        SELECT *
        FROM indicadores
        """,
        conn_dash,
    )

    hora_punta = pd.read_sql_query(
        """
        SELECT *
        FROM hora_punta
        """,
        conn_data,
    )

    conn_data.close()
    conn_dash.close()
    conn_insumos.close()

    # Imputar anio, mes y tipo de dia
    viajes['yr'] = pd.to_datetime(viajes.dia).dt.year
    viajes['mo'] = pd.to_datetime(viajes.dia).dt.month
    viajes['dow'] = pd.to_datetime(viajes.dia).dt.day_of_week

    viajes['tipo_dia'] = 'Dia habil'
    viajes.loc[viajes.dow >= 5, 'tipo_dia'] = 'Fin de semana'

    v_iter = viajes.groupby(['yr', 'mo', 'tipo_dia'],
                            as_index=False).size().iterrows()
    for _, i in v_iter:

        ym = f'{i.yr}-{str(i.mo).zfill(2)}'
        mo = str(i.mo).zfill(2)
        desc_dia = f'{i.yr}-{mo} ({i.tipo_dia})'
        desc_dia_titulo = f"Corrida de {meses[i.mo-1]} {i.yr} ({i.tipo_dia})"
        desc_dia_file = f'{i.yr}-{mo}({i.tipo_dia})'

        ind = indicadores[indicadores.dia.str[:7] == ym]
        ind = indicadores[indicadores.dia.str[:7] == ym]
        ind['ym'] = ym
        ind = ind.groupby(['ym', 'detalle'], as_index=False)[
            ['indicador', 'porcentaje']].mean()

        print(meses[int(ym[-2:])-1], ym[:4], desc_dia[8:])

        # Creo powerpoint
        prs = Presentation()
        prs.slide_height = Inches(13.5)
        prs.slide_width = Inches(24)

        # SLIDE 1 - Indicadores
        prs = slide_1(prs,
                      indicadores,
                      desc_dia=desc_dia,
                      tipo_dia=i.tipo_dia,
                      desc_dia_titulo=desc_dia_titulo,
                      )

        # SLIDE 2 -

        slide = get_new_slide(prs, desc_dia_titulo)

        punta_manana = hora_punta[(hora_punta.tipo_dia == i.tipo_dia) &
                                  (hora_punta.dia == f'{i.yr}-{mo}') &
                                  (hora_punta.detalle == 'Hora punta mañana')].indicador.astype(int).values[0]
        punta_mediodia = hora_punta[(hora_punta.tipo_dia == i.tipo_dia) &
                                    (hora_punta.dia == f'{i.yr}-{mo}') &
                                    (hora_punta.detalle == 'Hora punta mediodí­a')].indicador.astype(int).values[0]
        punta_tarde = hora_punta[(hora_punta.tipo_dia == i.tipo_dia) &
                                 (hora_punta.dia == f'{i.yr}-{mo}') &
                                 (hora_punta.detalle == 'Hora punta tarde')].indicador.astype(int).values[0]

        slide = pptx_text(prs=prs, slide=slide,  title=f'La hora punta mañana de viajes es a las {punta_manana} hs., la hora punta mediodía de viajes es a las {punta_mediodia} hs. y la hora punta tarde de viajes es a las {punta_tarde} hs.',
                          left=4, top=2.8, width=18, fontsize=20, bold=True)

        file_graph = os.path.join(
            "resultados",
            "png",
            f"{alias}{i.yr}-{mo}({i.tipo_dia})_viajes_x_hora.png")

        pptx_addpic(prs=prs,
                    slide=slide,
                    img_path=file_graph,
                    left=2.5,
                    top=4.5,
                    width=19)

        # SLIDE 3 -
        slide = get_new_slide(prs, desc_dia_titulo)

        hora_punta_modos = hora_punta[(hora_punta.tipo_dia == i.tipo_dia) &
                                      (hora_punta.dia == f'{i.yr}-{i.mo}') &
                                      (~hora_punta.detalle.isin(['Hora punta mañana', 'Hora punta mediodí­a', 'Hora punta tarde']))]
        hora_punta_modos = hora_punta_modos.reset_index().rename(columns={
            'index': 'orden'})
        hora_punta_modos = hora_punta_modos.sort_values(
            'orden', ascending=False).reset_index(drop=True)
        modos = []
        val = 0
        # for i in hora_punta_modos:
        for _ in range(0, len(hora_punta_modos), 3):
            det = hora_punta_modos.iloc[_].detalle
            det = det.split('(')[1][:-1]
            tarde = hora_punta_modos.iloc[_].indicador.astype(int)
            mediodia = hora_punta_modos.iloc[_+1].indicador.astype(int)
            manana = hora_punta_modos.iloc[_+2].indicador.astype(int)
            modos += [f'Viajes en {det}: hora punta mañana {manana} hs., hora punta mediodia {mediodia} hs., hora punta tarde {tarde} hs.']

        top_m = 2.5
        for modo in modos:
            slide = pptx_text(prs=prs, slide=slide,  title=modo,
                              left=4, top=top_m, width=18, fontsize=20, bold=True)
            top_m += .5

        file_graph = os.path.join(
            "resultados",
            "png",
            f"{alias}{i.yr}-{i.mo}({i.tipo_dia})_viajes_modo.png")

        pptx_addpic(prs=prs,
                    slide=slide,
                    img_path=file_graph,
                    left=2.5,
                    top=top_m+.5,
                    width=19)

        # SLIDE 4 -
        if i.tipo_dia == 'Dia habil':
            tip_dia = 'Hábil'
        else:
            tip_dia == 'Fin de semana'

        avg_dist = indicadores_dash[(indicadores_dash.desc_dia == f'{i.yr}-{mo}') &
                                    (indicadores_dash.tipo_dia == tip_dia) &
                                    (indicadores_dash.Indicador == 'Distancia de los viajes (promedio en kms)')].Valor.values[0].replace('\xa0', '')
        median_dist = indicadores_dash[(indicadores_dash.desc_dia == f'{i.yr}-{mo}') &
                                       (indicadores_dash.tipo_dia == tip_dia) &
                                       (indicadores_dash.Indicador == 'Distancia de los viajes (mediana en kms)')].Valor.values[0].replace('\xa0', '')

        slide = get_new_slide(prs, desc_dia_titulo)

        slide = pptx_text(prs=prs, slide=slide,  title=f'La distancia promedio de los viajes es de {avg_dist} y la distancia mediana es de {median_dist}',
                          left=4, top=2.8, width=18, fontsize=20, bold=True)

        file_graph = os.path.join(
            "resultados",
            "png",
            f"{alias}{i.yr}-{mo}({i.tipo_dia})_viajes_dist.png")

        pptx_addpic(prs=prs,
                    slide=slide,
                    img_path=file_graph,
                    left=4,
                    top=4,
                    width=15)

        # SLIDE 5 -

        slide = get_new_slide(prs, desc_dia_titulo)

        file_graph = os.path.join(
            "resultados",
            "png",
            f"{alias}{i.yr}-{mo}({i.tipo_dia})_{geo_files[0][1]}_lineas_deseo.png")

        pptx_addpic(prs=prs,
                    slide=slide,
                    img_path=file_graph,
                    left=1,
                    top=2.5,
                    width=9)

        file_graph = os.path.join(
            "resultados",
            "png",
            f"{alias}{i.yr}-{mo}({i.tipo_dia})_{geo_files[0][1]}_matrizod.png")

        pptx_addpic(prs=prs,
                    slide=slide,
                    img_path=file_graph,
                    left=14,
                    top=2.5,
                    width=7)

        # SLIDE 6 -
        conn_data = iniciar_conexion_db(tipo='data')

        q = """
            select id_linea
            from etapas e
            where modo = 'autobus'
            group by id_linea
            order by sum(factor_expansion_linea) DESC
            limit 1;
        """
        id_linea = pd.read_sql(q, conn_data)
        id_linea = id_linea.id_linea.item()

        if i.tipo_dia == 'Dia habil':
            tdia = 'weekday'
        else:
            tdia = 'weekend'

        # basic kpi plot
        slide = get_new_slide(prs, desc_dia_titulo)

        file_graph = os.path.join(
            "resultados",
            "png",
            f"{alias}_kpi_basicos_id_linea_{id_linea}_{tdia}.png")

        if os.path.isfile(file_graph):

            pptx_addpic(prs=prs,
                        slide=slide,
                        img_path=file_graph,
                        left=1,
                        top=2.5,
                        width=10)
            # query demand data
            q = f"""
                select *
                from basic_kpi_by_line_day
                where id_linea = {id_linea}
                and dia = '{tdia}';
            """
            demand_data = pd.read_sql(q, conn_data)
            pax = int(demand_data.pax.item())
            veh = int(demand_data.veh.item())
            dmt = int(demand_data.dmt.item())
            speed = int(demand_data.speed_kmh.item())

            s = f"La linea id {id_linea} tiene en {i.tipo_dia} "
            s = s + f"una demanda de {pax} pasajeros, con "
            s = s + \
                f"{veh} vehiculos dia, una distancia media de {dmt} metros "
            s = s + \
                f"y una velocidad comercial promedio de {speed} kmh"

            # create plot and text
            slide = pptx_text(
                prs=prs, slide=slide,
                title=s,
                left=1, top=11, width=10, fontsize=20, bold=True)

        else:
            print("No existe el archivo", file_graph)

        # section load plot
        file_graph = os.path.join(
            "resultados",
            "png",
            f"{alias}_{tdia}_segmentos_id_linea_{id_linea}_prop_etapas_ 7-10 hrs.png")

        if os.path.isfile(file_graph):

            pptx_addpic(prs=prs,
                        slide=slide,
                        img_path=file_graph,
                        left=13,
                        top=2.5,
                        width=10)

            q = f"""
                select round(section_id * 100,1) as section_id,
                       sentido, cantidad_etapas
                from ocupacion_por_linea_tramo
                where id_linea = {id_linea}
                and day_type = '{tdia}'
                and n_sections = 10
                and hora_min = 7
                and hora_max = 10
                order by cantidad_etapas desc
                limit 1;
            """
            load_data = pd.read_sql(q, conn_data)
            section_id = load_data.section_id.item()
            sentido = load_data.sentido.item()
            cantidad_etapas = int(load_data.cantidad_etapas.item())

            s = f"Y su tramo mas cargado es en el sentido {sentido} "
            s = s + \
                f" en el tramo equivalente al {section_id} % de su recorrido "
            s = s + \
                f" con {cantidad_etapas} etapas."

            # create plot and text
            slide = pptx_text(
                prs=prs, slide=slide,
                title=s,
                left=13, top=11, width=10, fontsize=20, bold=True)

        else:
            print("No existe el archivo", file_graph)

        try:
            file_pptx = os.path.join(
                "resultados", "ppts", f"{alias}{desc_dia_file}.pptx")
            print('')
            prs.save(file_pptx)
            print(file_pptx)
            print('')
        except:
            print('')
            print('No se pudo guardar el archivo', file_pptx)
